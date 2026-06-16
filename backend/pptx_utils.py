"""
Utilities for copying slides between presentations and merging pptx files.
"""
from pptx import Presentation
from pptx.util import Pt
from lxml import etree
import copy
import io


def get_slide_texts(pptx_bytes: bytes) -> list[dict]:
    """Extract text content from each slide in a pptx file."""
    prs = Presentation(io.BytesIO(pptx_bytes))
    results = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text[:120])
        results.append({"index": i, "texts": texts[:6]})
    return results


def extract_slides_as_pptx(source_bytes: bytes, slide_indices: list[int]) -> bytes:
    """
    Extract specific slides from a pptx and return a new pptx containing only those slides.
    slide_indices: 0-based list of indices to extract, in order.
    """
    source = Presentation(io.BytesIO(source_bytes))
    target = Presentation()
    target.slide_width = source.slide_width
    target.slide_height = source.slide_height

    for idx in slide_indices:
        if 0 <= idx < len(source.slides):
            _copy_slide(target, source.slides[idx])

    buf = io.BytesIO()
    target.save(buf)
    return buf.getvalue()


def merge_pptx_list(pptx_bytes_list: list[bytes]) -> bytes:
    """
    Merge multiple pptx files (each as bytes) into one presentation.
    Uses the first file's slide dimensions and master theme.
    """
    if not pptx_bytes_list:
        prs = Presentation()
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    merged = Presentation(io.BytesIO(pptx_bytes_list[0]))

    for pptx_bytes in pptx_bytes_list[1:]:
        source = Presentation(io.BytesIO(pptx_bytes))
        for slide in source.slides:
            _copy_slide(merged, slide)

    buf = io.BytesIO()
    merged.save(buf)
    return buf.getvalue()


def _copy_slide(target_prs: Presentation, source_slide) -> None:
    """
    Copy a slide from source into target_prs, preserving images and all content.
    """
    blank_layout = target_prs.slide_layouts[-1]
    new_slide = target_prs.slides.add_slide(blank_layout)

    # Build rId mapping: source rId -> new rId in the destination slide
    rId_map: dict[str, str] = {}
    for rId, rel in source_slide.part.rels.items():
        if not rel.is_external:
            try:
                new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
                rId_map[rId] = new_rId
            except Exception:
                pass

    # Deep copy source slide XML and update relationship ID references
    src_elem = copy.deepcopy(source_slide._element)
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    for elem in src_elem.iter():
        for attr in list(elem.attrib.keys()):
            if r_ns in str(attr):
                old_val = elem.get(attr)
                if old_val and old_val in rId_map:
                    elem.set(attr, rId_map[old_val])

    # Replace new slide content with source content
    new_slide_elem = new_slide._element
    for child in list(new_slide_elem):
        new_slide_elem.remove(child)
    for child in src_elem:
        new_slide_elem.append(copy.deepcopy(child))
    for attr, val in src_elem.attrib.items():
        new_slide_elem.set(attr, val)


def count_slides(pptx_bytes: bytes) -> int:
    """Count the number of slides in a pptx file."""
    prs = Presentation(io.BytesIO(pptx_bytes))
    return len(prs.slides)
